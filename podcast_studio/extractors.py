"""Extract plain text from uploaded documents (PDF, Word, PowerPoint, Excel, text)."""

import io
from pathlib import Path


def extract_text(filename: str, data: bytes) -> str:
    """Return the text content of an uploaded file, or an error note if unreadable."""
    suffix = Path(filename).suffix.lower()
    try:
        if suffix == ".pdf":
            text = _from_pdf(data)
        elif suffix == ".docx":
            text = _from_docx(data)
        elif suffix == ".pptx":
            text = _from_pptx(data)
        elif suffix in (".xlsx", ".xlsm"):
            text = _from_xlsx(data)
        else:
            # Fall back to plain-text decoding (.txt, .md, .csv, .json, ...)
            text = data.decode("utf-8", errors="replace")
        # CRM/Salesforce exports embed literal carriage-return escapes.
        return text.replace("_x000D_", "")
    except Exception as exc:  # surface per-file failures without killing the batch
        return f"[Could not read {filename}: {exc}]"


def _from_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _from_docx(data: bytes) -> str:
    import docx

    document = docx.Document(io.BytesIO(data))
    parts = [p.text for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _from_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append("\t".join(cells))
    wb.close()
    return "\n".join(parts)
